"""Извлечение «частей» ввода из Telegram-сообщения: текст, PDF, фото, голос."""
import asyncio
import io
import os
import shutil
import subprocess
import tempfile

from . import config


def parse_deck_hint(text: str | None) -> str | None:
    """Вытащить колоду из строк 'Колода:' / 'Подколода:'. None, если не указана."""
    if not text:
        return None
    deck = None
    subdeck = None
    for line in text.splitlines():
        line = line.strip()
        low = line.lower()
        if low.startswith("колода:") or low.startswith("deck:"):
            deck = line.split(":", 1)[1].strip()
        elif low.startswith("подколода:") or low.startswith("subdeck:"):
            subdeck = line.split(":", 1)[1].strip()
    if deck and subdeck:
        return f"{deck}::{subdeck}"
    return deck


async def _download(context, file_id: str) -> bytes:
    f = await context.bot.get_file(file_id)
    return bytes(await f.download_as_bytearray())


def _pdf_to_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages).strip()
    except Exception:
        return ""


def _docx_fallback(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    out = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out).strip()


def _pptx_fallback(data: bytes) -> str:
    from pptx import Presentation
    out = []
    for slide in Presentation(io.BytesIO(data)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
    return "\n".join(out).strip()


def _doc_via_textutil(data: bytes) -> str:
    """Старый .doc → текст через встроенный macOS textutil (нет на Linux)."""
    if not shutil.which("textutil"):
        return ""
    with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as f:
        f.write(data)
        path = f.name
    try:
        out = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", path],
            capture_output=True, timeout=30,
        )
        return out.stdout.decode("utf-8", errors="replace").strip()
    except Exception:
        return ""
    finally:
        os.unlink(path)


def _xls_to_text(data: bytes) -> str:
    """Старый .xls → текст через xlrd (чистый Python, работает везде)."""
    import xlrd
    out = []
    book = xlrd.open_workbook(file_contents=data)
    for sheet in book.sheets():
        out.append(f"# {sheet.name}")
        for r in range(sheet.nrows):
            row = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
            if any(row):
                out.append(" | ".join(row))
    return "\n".join(out).strip()


def _office_to_markdown(data: bytes, filename: str) -> str:
    """Word / PowerPoint / Excel → текст. markitdown, с фолбэком на прямые библиотеки."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        from markitdown import MarkItDown
        res = MarkItDown().convert_stream(io.BytesIO(data), file_extension=ext)
        text = (res.text_content or "").strip()
    except Exception:
        text = ""
    if not text and ext == ".docx":
        text = _docx_fallback(data)
    elif not text and ext == ".pptx":
        text = _pptx_fallback(data)
    return text


_OFFICE_EXT = (".docx", ".pptx", ".xlsx")
_OFFICE_MIME = ("officedocument", "msword", "ms-powerpoint", "ms-excel")
_OFFICE_LEGACY_EXT = (".doc", ".ppt", ".xls")
_OFFICE_LEGACY_MIME = ("application/msword", "application/vnd.ms-powerpoint", "application/vnd.ms-excel")


def _transcribe(data: bytes, filename: str) -> str:
    from groq import Groq
    client = Groq(api_key=config.GROQ_API_KEY)
    tr = client.audio.transcriptions.create(
        file=(filename, data),
        model=config.GROQ_WHISPER_MODEL,
    )
    return (tr.text or "").strip()


async def extract(update, context) -> tuple[list[dict], str | None, str | None]:
    """Вернуть (parts, deck_hint, note).

    note — необязательное предупреждение для пользователя (например, нет ключа
    для расшифровки голоса).
    """
    msg = update.message
    parts: list[dict] = []
    note: str | None = None

    text = msg.text or msg.caption
    deck_hint = parse_deck_hint(text)

    # --- Фото ---
    if msg.photo:
        data = await _download(context, msg.photo[-1].file_id)
        parts.append({"type": "image", "mime": "image/jpeg", "data": data})

    # --- Документ (PDF / txt / картинка) ---
    if msg.document:
        doc = msg.document
        mime = (doc.mime_type or "").lower()
        name = (doc.file_name or "").lower()
        data = await _download(context, doc.file_id)
        if mime == "application/pdf" or name.endswith(".pdf"):
            extracted = _pdf_to_text(data)
            if len(extracted) > 200:
                parts.append({"type": "text", "text": extracted})
            else:
                parts.append({"type": "pdf", "data": data})  # вероятно скан → мультимодал
        elif mime.startswith("image/"):
            parts.append({"type": "image", "mime": mime, "data": data})
        elif name.endswith(_OFFICE_EXT) or any(k in mime for k in _OFFICE_MIME):
            is_doc = name.endswith(".doc") or mime == "application/msword"
            is_xls = name.endswith(".xls") or mime == "application/vnd.ms-excel"
            is_legacy = name.endswith(_OFFICE_LEGACY_EXT) or mime in _OFFICE_LEGACY_MIME
            if is_doc:
                txt = await asyncio.to_thread(_doc_via_textutil, data)
                if txt:
                    parts.append({"type": "text", "text": txt})
                else:
                    note = ("⚠️ Старый .doc читается только локальным ботом на Mac. "
                            "Сохрани как .docx и пришли снова.")
            elif is_xls:
                try:
                    txt = await asyncio.to_thread(_xls_to_text, data)
                except Exception as e:
                    txt = ""
                    note = f"⚠️ Не смог прочитать .xls: {e}"
                if txt:
                    parts.append({"type": "text", "text": txt})
                elif not note:
                    note = "⚠️ В таблице .xls не найден текст."
            elif is_legacy:
                note = ("⚠️ Старый .ppt не поддерживается напрямую. "
                        "Сохрани как .pptx и пришли снова.")
            else:
                try:
                    md = await asyncio.to_thread(_office_to_markdown, data, name)
                except Exception as e:
                    md = ""
                    note = f"⚠️ Не смог обработать «{doc.file_name}»: {e}"
                if md:
                    parts.append({"type": "text", "text": md})
                elif not note:
                    note = f"⚠️ В «{doc.file_name}» не найден текст (возможно, внутри только картинки/сканы)."
        elif mime.startswith("text/") or name.endswith((".txt", ".md")):
            parts.append({"type": "text", "text": data.decode("utf-8", errors="replace")})
        else:
            note = f"⚠️ Формат {mime or name or 'файла'} не поддерживается — пропущен."

    # --- Голос / аудио ---
    if msg.voice or msg.audio:
        media = msg.voice or msg.audio
        if not config.GROQ_API_KEY:
            note = "⚠️ Для расшифровки голоса нужен GROQ_API_KEY. Сообщение пропущено."
        else:
            data = await _download(context, media.file_id)
            fname = "voice.ogg" if msg.voice else (getattr(media, "file_name", None) or "audio.mp3")
            try:
                transcript = await asyncio.to_thread(_transcribe, data, fname)
                if transcript:
                    parts.append({"type": "text", "text": transcript})
                    note = f"🎤 Распознал: «{transcript[:300]}»"
                else:
                    note = "⚠️ Не удалось распознать речь в сообщении."
            except Exception as e:
                note = f"⚠️ Ошибка расшифровки голоса: {e}"

    # --- Текст / подпись ---
    if text:
        parts.append({"type": "text", "text": text})

    return parts, deck_hint, note
